"""Tests for the MCP server tool functions."""

from unittest.mock import MagicMock, patch

import pytest

from outpost.mcp_server import (
    auth_status,
    cal_add,
    cal_delete,
    cal_list,
    cal_next,
    cal_today,
    cal_update,
    contact_list,
    contact_search,
    mail_attachments,
    mail_delete,
    mail_download_attachment,
    mail_list,
    mail_read,
    mail_reply,
    mail_search,
    mail_send,
    task_add,
    task_complete,
    task_delete,
    task_list,
    task_lists,
    task_lists_create,
    task_lists_delete,
    task_update,
    teams_channels,
    teams_download,
    teams_files,
    teams_list,
    teams_messages,
    teams_send,
    teams_upload,
    teams_workspace_extract,
    teams_workspace_list,
    teams_workspace_read,
    teams_workspace_write,
)


# ── Auth error path ─────────────────────────────────────────────────────────


@patch("outpost.mcp_server.get_token", return_value=None)
def test_unauthenticated_raises_runtime_error(mock_token):
    """Tools raise RuntimeError with helpful message when not authenticated."""
    with pytest.raises(RuntimeError, match="Not authenticated"):
        task_list()


# ── task_add ────────────────────────────────────────────────────────────────


@patch("outpost.mcp_server.get_token", return_value="fake-token")
@patch("outpost.api.tasks.add_task")
@patch("outpost.mcp_server.GraphClient")
def test_task_add_basic(mock_client_cls, mock_add_task, mock_token):
    mock_add_task.return_value = {"id": "1", "title": "Buy milk"}
    result = task_add(title="Buy milk")
    mock_add_task.assert_called_once_with(
        mock_client_cls.return_value, "Buy milk", due_date=None, list_name=None, priority="normal"
    )
    assert result == {"id": "1", "title": "Buy milk"}


@patch("outpost.mcp_server.get_token", return_value="fake-token")
@patch("outpost.api.tasks.add_task")
@patch("outpost.mcp_server.GraphClient")
def test_task_add_with_due_and_list(mock_client_cls, mock_add_task, mock_token):
    mock_add_task.return_value = {"id": "2", "title": "Report"}
    result = task_add(title="Report", due="2026-03-15", list_name="Work", priority="high")
    mock_add_task.assert_called_once_with(
        mock_client_cls.return_value, "Report", due_date="2026-03-15", list_name="Work", priority="high"
    )
    assert result["title"] == "Report"


# ── task_list ───────────────────────────────────────────────────────────────


@patch("outpost.mcp_server.get_token", return_value="fake-token")
@patch("outpost.api.tasks.list_tasks")
@patch("outpost.mcp_server.GraphClient")
def test_task_list_no_filter(mock_client_cls, mock_list_tasks, mock_token):
    mock_list_tasks.return_value = [{"id": "1", "title": "Test"}]
    result = task_list()
    mock_list_tasks.assert_called_once_with(
        mock_client_cls.return_value, due_filter=None, list_name=None
    )
    assert len(result) == 1


@patch("outpost.mcp_server.get_token", return_value="fake-token")
@patch("outpost.api.tasks.list_tasks")
@patch("outpost.mcp_server.GraphClient")
def test_task_list_with_due(mock_client_cls, mock_list_tasks, mock_token):
    mock_list_tasks.return_value = []
    result = task_list(due="2026-03-10")
    mock_list_tasks.assert_called_once_with(
        mock_client_cls.return_value, due_filter="2026-03-10", list_name=None
    )
    assert result == []


# ── cal_add ─────────────────────────────────────────────────────────────────


@patch("outpost.mcp_server.get_token", return_value="fake-token")
@patch("outpost.api.calendar.create_event")
@patch("outpost.mcp_server.GraphClient")
def test_cal_add_basic(mock_client_cls, mock_create, mock_token):
    mock_create.return_value = {"id": "e1", "subject": "Meeting"}
    result = cal_add(title="Meeting", start="2026-03-15T09:00")
    assert mock_create.call_count == 1
    call_args = mock_create.call_args
    assert call_args[0][1] == "Meeting"  # title
    assert result["subject"] == "Meeting"


@patch("outpost.mcp_server.get_token", return_value="fake-token")
@patch("outpost.api.calendar.create_event")
@patch("outpost.mcp_server.GraphClient")
def test_cal_add_with_duration(mock_client_cls, mock_create, mock_token):
    mock_create.return_value = {"id": "e2", "subject": "Standup"}
    result = cal_add(title="Standup", start="2026-03-15T09:00", duration=15)
    call_args = mock_create.call_args
    assert call_args[1]["duration"] == 15


# ── cal_today ───────────────────────────────────────────────────────────────


@patch("outpost.mcp_server.get_token", return_value="fake-token")
@patch("outpost.api.calendar.get_today_events")
@patch("outpost.mcp_server.GraphClient")
def test_cal_today(mock_client_cls, mock_today, mock_token):
    mock_today.return_value = [{"subject": "Lunch"}]
    result = cal_today()
    mock_today.assert_called_once_with(mock_client_cls.return_value)
    assert result == [{"subject": "Lunch"}]


# ── cal_list ────────────────────────────────────────────────────────────────


@patch("outpost.mcp_server.get_token", return_value="fake-token")
@patch("outpost.api.calendar.get_week_events")
@patch("outpost.mcp_server.GraphClient")
def test_cal_list_default_week(mock_client_cls, mock_week, mock_token):
    mock_week.return_value = [{"subject": "Weekly"}]
    result = cal_list()
    mock_week.assert_called_once_with(mock_client_cls.return_value)
    assert len(result) == 1


@patch("outpost.mcp_server.get_token", return_value="fake-token")
@patch("outpost.api.calendar.get_calendar_view")
@patch("outpost.mcp_server.GraphClient")
def test_cal_list_specific_date(mock_client_cls, mock_view, mock_token):
    mock_view.return_value = [{"subject": "Birthday"}]
    result = cal_list(date="2026-03-15")
    assert mock_view.call_count == 1
    assert result == [{"subject": "Birthday"}]


# ── auth_status ─────────────────────────────────────────────────────────────


# ── task_update ─────────────────────────────────────────────────────────────


@patch("outpost.mcp_server.get_token", return_value="fake-token")
@patch("outpost.api.tasks.update_task")
@patch("outpost.mcp_server.GraphClient")
def test_task_update(mock_client_cls, mock_update, mock_token):
    mock_update.return_value = {"id": "1", "title": "Updated"}
    result = task_update(task_id="1", title="Updated")
    mock_update.assert_called_once()
    assert result["title"] == "Updated"


# ── task_complete ──────────────────────────────────────────────────────────


@patch("outpost.mcp_server.get_token", return_value="fake-token")
@patch("outpost.api.tasks.complete_task")
@patch("outpost.mcp_server.GraphClient")
def test_task_complete(mock_client_cls, mock_complete, mock_token):
    mock_complete.return_value = {"id": "1", "status": "completed"}
    result = task_complete(task_id="1")
    mock_complete.assert_called_once()
    assert result["status"] == "completed"


# ── task_delete ────────────────────────────────────────────────────────────


@patch("outpost.mcp_server.get_token", return_value="fake-token")
@patch("outpost.api.tasks.delete_task")
@patch("outpost.mcp_server.GraphClient")
def test_task_delete(mock_client_cls, mock_del, mock_token):
    mock_del.return_value = {}
    result = task_delete(task_id="1")
    mock_del.assert_called_once()
    assert result == {}


# ── cal_update ─────────────────────────────────────────────────────────────


@patch("outpost.mcp_server.get_token", return_value="fake-token")
@patch("outpost.api.calendar.update_event")
@patch("outpost.mcp_server.GraphClient")
def test_cal_update(mock_client_cls, mock_update, mock_token):
    mock_update.return_value = {"id": "e1", "subject": "Updated"}
    result = cal_update(event_id="e1", title="Updated")
    assert mock_update.call_count == 1
    assert result["subject"] == "Updated"


# ── cal_delete ─────────────────────────────────────────────────────────────


@patch("outpost.mcp_server.get_token", return_value="fake-token")
@patch("outpost.api.calendar.delete_event")
@patch("outpost.mcp_server.GraphClient")
def test_cal_delete(mock_client_cls, mock_del, mock_token):
    mock_del.return_value = {}
    result = cal_delete(event_id="e1")
    mock_del.assert_called_once()
    assert result == {}


# ── mail_list ──────────────────────────────────────────────────────────────


@patch("outpost.mcp_server.get_token", return_value="fake-token")
@patch("outpost.api.mail.list_messages")
@patch("outpost.mcp_server.GraphClient")
def test_mail_list(mock_client_cls, mock_list, mock_token):
    mock_list.return_value = [{"id": "m1", "subject": "Hello"}]
    result = mail_list()
    mock_list.assert_called_once_with(mock_client_cls.return_value, folder="inbox", top=25, unread=False)
    assert len(result) == 1


# ── mail_read ──────────────────────────────────────────────────────────────


@patch("outpost.mcp_server.get_token", return_value="fake-token")
@patch("outpost.api.mail.get_message")
@patch("outpost.mcp_server.GraphClient")
def test_mail_read(mock_client_cls, mock_get, mock_token):
    mock_get.return_value = {"id": "m1", "subject": "Hello", "body": {"content": "Hi"}}
    result = mail_read(message_id="m1")
    mock_get.assert_called_once_with(mock_client_cls.return_value, "m1")
    assert result["subject"] == "Hello"


# ── mail_send ──────────────────────────────────────────────────────────────


@patch("outpost.mcp_server.get_token", return_value="fake-token")
@patch("outpost.api.mail.send_message")
@patch("outpost.mcp_server.GraphClient")
def test_mail_send(mock_client_cls, mock_send, mock_token):
    mock_send.return_value = {}
    result = mail_send(to="bob@example.com", subject="Test", body="Hi")
    mock_send.assert_called_once_with(
        mock_client_cls.return_value, ["bob@example.com"], "Test", "Hi", cc=None
    )
    assert result == {}


@patch("outpost.mcp_server.get_token", return_value="fake-token")
@patch("outpost.api.mail.send_message")
@patch("outpost.mcp_server.GraphClient")
def test_mail_send_with_cc(mock_client_cls, mock_send, mock_token):
    mock_send.return_value = {}
    mail_send(to="bob@example.com", subject="Test", body="Hi", cc="carol@example.com")
    mock_send.assert_called_once_with(
        mock_client_cls.return_value, ["bob@example.com"], "Test", "Hi", cc=["carol@example.com"]
    )


# ── mail_reply ─────────────────────────────────────────────────────────────


@patch("outpost.mcp_server.get_token", return_value="fake-token")
@patch("outpost.api.mail.reply_message")
@patch("outpost.mcp_server.GraphClient")
def test_mail_reply(mock_client_cls, mock_reply, mock_token):
    mock_reply.return_value = {}
    result = mail_reply(message_id="m1", body="Thanks!")
    mock_reply.assert_called_once_with(mock_client_cls.return_value, "m1", "Thanks!")
    assert result == {}


# ── mail_delete ────────────────────────────────────────────────────────────


@patch("outpost.mcp_server.get_token", return_value="fake-token")
@patch("outpost.api.mail.delete_message")
@patch("outpost.mcp_server.GraphClient")
def test_mail_delete(mock_client_cls, mock_del, mock_token):
    mock_del.return_value = {}
    result = mail_delete(message_id="m1")
    mock_del.assert_called_once_with(mock_client_cls.return_value, "m1")
    assert result == {}


# ── contact_list ───────────────────────────────────────────────────────────


@patch("outpost.mcp_server.get_token", return_value="fake-token")
@patch("outpost.api.contacts.list_contacts")
@patch("outpost.mcp_server.GraphClient")
def test_contact_list(mock_client_cls, mock_list, mock_token):
    mock_list.return_value = [{"displayName": "Alice"}]
    result = contact_list()
    mock_list.assert_called_once_with(mock_client_cls.return_value, top=50)
    assert len(result) == 1


# ── contact_search ─────────────────────────────────────────────────────────


@patch("outpost.mcp_server.get_token", return_value="fake-token")
@patch("outpost.api.contacts.search_contacts")
@patch("outpost.mcp_server.GraphClient")
def test_contact_search(mock_client_cls, mock_search, mock_token):
    mock_search.return_value = [{"displayName": "Alice"}]
    result = contact_search(query="Alice")
    mock_search.assert_called_once_with(mock_client_cls.return_value, "Alice", top=25)
    assert len(result) == 1


# ── auth_status ─────────────────────────────────────────────────────────────


@patch("outpost.mcp_server.get_auth_status")
def test_auth_status_logged_in(mock_status):
    mock_status.return_value = {"logged_in": True, "username": "user@example.com"}
    result = auth_status()
    assert result["logged_in"] is True
    assert result["username"] == "user@example.com"


@patch("outpost.mcp_server.get_auth_status")
def test_auth_status_not_logged_in(mock_status):
    mock_status.return_value = {"logged_in": False, "username": None}
    result = auth_status()
    assert result["logged_in"] is False


# ── cal_next ──────────────────────────────────────────────────────────────


@patch("outpost.mcp_server.get_token", return_value="fake-token")
@patch("outpost.api.calendar.get_next_events")
@patch("outpost.mcp_server.GraphClient")
def test_cal_next(mock_client_cls, mock_next, mock_token):
    mock_next.return_value = [{"subject": "Next meeting"}]
    result = cal_next(count=1)
    mock_next.assert_called_once_with(mock_client_cls.return_value, count=1)
    assert len(result) == 1


# ── mail_search ───────────────────────────────────────────────────────────


@patch("outpost.mcp_server.get_token", return_value="fake-token")
@patch("outpost.api.mail.search_messages")
@patch("outpost.mcp_server.GraphClient")
def test_mail_search(mock_client_cls, mock_search, mock_token):
    mock_search.return_value = [{"subject": "Budget report"}]
    result = mail_search(query="budget")
    mock_search.assert_called_once_with(mock_client_cls.return_value, "budget", top=25)
    assert len(result) == 1


# ── mail_attachments ─────────────────────────────────────────────────────


@patch("outpost.mcp_server.get_token", return_value="fake-token")
@patch("outpost.api.mail.list_attachments")
@patch("outpost.mcp_server.GraphClient")
def test_mail_attachments(mock_client_cls, mock_list_att, mock_token):
    mock_list_att.return_value = [{"id": "att-1", "name": "file.pdf"}]
    result = mail_attachments(message_id="msg-1")
    mock_list_att.assert_called_once_with(mock_client_cls.return_value, "msg-1")
    assert len(result) == 1


# ── mail_download_attachment ──────────────────────────────────────────────


@patch("outpost.mcp_server.get_token", return_value="fake-token")
@patch("outpost.api.mail.download_attachment")
@patch("outpost.mcp_server.GraphClient")
def test_mail_download_attachment(mock_client_cls, mock_download, mock_token):
    mock_download.return_value = ("test.txt", b"hello")
    result = mail_download_attachment(message_id="msg-1", attachment_id="att-1")
    assert result["filename"] == "test.txt"
    assert "content_base64" in result


# ── task_lists ────────────────────────────────────────────────────────────


@patch("outpost.mcp_server.get_token", return_value="fake-token")
@patch("outpost.api.tasks.get_task_lists")
@patch("outpost.mcp_server.GraphClient")
def test_task_lists(mock_client_cls, mock_get_lists, mock_token):
    mock_get_lists.return_value = [{"id": "l1", "displayName": "Tasks"}]
    result = task_lists()
    mock_get_lists.assert_called_once_with(mock_client_cls.return_value)
    assert len(result) == 1


# ── task_lists_create ─────────────────────────────────────────────────────


@patch("outpost.mcp_server.get_token", return_value="fake-token")
@patch("outpost.api.tasks.create_task_list")
@patch("outpost.mcp_server.GraphClient")
def test_task_lists_create(mock_client_cls, mock_create, mock_token):
    mock_create.return_value = {"id": "new", "displayName": "Shopping"}
    result = task_lists_create(name="Shopping")
    mock_create.assert_called_once_with(mock_client_cls.return_value, "Shopping")
    assert result["displayName"] == "Shopping"


# ── task_lists_delete ─────────────────────────────────────────────────────


@patch("outpost.mcp_server.get_token", return_value="fake-token")
@patch("outpost.api.tasks.delete_task_list")
@patch("outpost.mcp_server.GraphClient")
def test_task_lists_delete(mock_client_cls, mock_del, mock_token):
    mock_del.return_value = {}
    result = task_lists_delete(list_id="list-1")
    mock_del.assert_called_once_with(mock_client_cls.return_value, "list-1")
    assert result == {}


# ── cal_add with attendees ────────────────────────────────────────────────


@patch("outpost.mcp_server.get_token", return_value="fake-token")
@patch("outpost.api.calendar.create_event")
@patch("outpost.mcp_server.GraphClient")
def test_cal_add_with_attendees(mock_client_cls, mock_create, mock_token):
    mock_create.return_value = {"id": "e1", "subject": "Meeting"}
    result = cal_add(title="Meeting", start="2026-03-15T09:00", attendees="alice@example.com,bob@example.com")
    call_args = mock_create.call_args
    assert call_args[1]["attendees"] == ["alice@example.com", "bob@example.com"]


# ── cal_add / cal_update with show_as ────────────────────────────────────


@patch("outpost.mcp_server.get_token", return_value="fake-token")
@patch("outpost.api.calendar.create_event")
@patch("outpost.mcp_server.GraphClient")
def test_cal_add_with_show_as(mock_client_cls, mock_create, mock_token):
    mock_create.return_value = {"id": "e1", "subject": "PTO", "showAs": "oof"}
    result = cal_add(title="PTO", start="2026-03-15T09:00", show_as="oof")
    call_args = mock_create.call_args
    assert call_args[1]["show_as"] == "oof"
    assert result["showAs"] == "oof"


@patch("outpost.mcp_server.get_token", return_value="fake-token")
@patch("outpost.api.calendar.update_event")
@patch("outpost.mcp_server.GraphClient")
def test_cal_update_with_show_as(mock_client_cls, mock_update, mock_token):
    mock_update.return_value = {"id": "e1", "showAs": "free"}
    result = cal_update(event_id="e1", show_as="free")
    call_args = mock_update.call_args
    assert call_args[1]["show_as"] == "free"
    assert result["showAs"] == "free"


# ── cal_add / cal_update with timezone ───────────────────────────────────


@patch("outpost.mcp_server.load_config", return_value={"timezone": "America/New_York"})
@patch("outpost.mcp_server.get_token", return_value="fake-token")
@patch("outpost.api.calendar.create_event")
@patch("outpost.mcp_server.GraphClient")
def test_cal_add_uses_config_timezone(mock_client_cls, mock_create, mock_token, mock_config):
    mock_create.return_value = {"id": "e1", "subject": "Meeting"}
    cal_add(title="Meeting", start="2026-03-15T09:00")
    call_args = mock_create.call_args
    assert call_args[1]["timezone"] == "America/New_York"


@patch("outpost.mcp_server.load_config", return_value={"timezone": "America/New_York"})
@patch("outpost.mcp_server.get_token", return_value="fake-token")
@patch("outpost.api.calendar.create_event")
@patch("outpost.mcp_server.GraphClient")
def test_cal_add_timezone_override(mock_client_cls, mock_create, mock_token, mock_config):
    mock_create.return_value = {"id": "e1", "subject": "Meeting"}
    cal_add(title="Meeting", start="2026-03-15T09:00", timezone="Europe/London")
    call_args = mock_create.call_args
    assert call_args[1]["timezone"] == "Europe/London"


@patch("outpost.mcp_server.load_config", return_value={"timezone": "America/Chicago"})
@patch("outpost.mcp_server.get_token", return_value="fake-token")
@patch("outpost.api.calendar.update_event")
@patch("outpost.mcp_server.GraphClient")
def test_cal_update_uses_config_timezone(mock_client_cls, mock_update, mock_token, mock_config):
    mock_update.return_value = {"id": "e1"}
    cal_update(event_id="e1", start="2026-03-20T10:00")
    call_args = mock_update.call_args
    assert call_args[1]["timezone"] == "America/Chicago"


# ── mail_list with unread ────────────────────────────────────────────────


@patch("outpost.mcp_server.get_token", return_value="fake-token")
@patch("outpost.api.mail.list_messages")
@patch("outpost.mcp_server.GraphClient")
def test_mail_list_unread(mock_client_cls, mock_list, mock_token):
    mock_list.return_value = []
    result = mail_list(unread=True)
    mock_list.assert_called_once_with(mock_client_cls.return_value, folder="inbox", top=25, unread=True)
    assert result == []


# ── Teams tools ──────────────────────────────────────────────────────────


@patch("outpost.config.is_feature_enabled", return_value=False)
def test_teams_list_disabled(mock_feature):
    with pytest.raises(RuntimeError, match="not enabled"):
        teams_list()


@patch("outpost.config.is_feature_enabled", return_value=True)
@patch("outpost.mcp_server.get_token", return_value="fake-token")
@patch("outpost.api.teams.list_teams")
@patch("outpost.mcp_server.GraphClient")
def test_teams_list(mock_client_cls, mock_list, mock_token, mock_feature):
    mock_list.return_value = [{"id": "t1", "displayName": "Eng"}]
    result = teams_list()
    mock_list.assert_called_once_with(mock_client_cls.return_value)
    assert len(result) == 1


@patch("outpost.config.is_feature_enabled", return_value=True)
@patch("outpost.mcp_server.get_token", return_value="fake-token")
@patch("outpost.api.teams.list_channels")
@patch("outpost.mcp_server.GraphClient")
def test_teams_channels(mock_client_cls, mock_list, mock_token, mock_feature):
    mock_list.return_value = [{"id": "ch-1", "displayName": "General"}]
    result = teams_channels(team_id="t1")
    mock_list.assert_called_once_with(mock_client_cls.return_value, "t1")
    assert len(result) == 1


@patch("outpost.config.is_feature_enabled", return_value=True)
@patch("outpost.mcp_server.get_token", return_value="fake-token")
@patch("outpost.api.teams.list_messages")
@patch("outpost.mcp_server.GraphClient")
def test_teams_messages(mock_client_cls, mock_list, mock_token, mock_feature):
    mock_list.return_value = [{"id": "m1", "body": {"content": "Hi"}}]
    result = teams_messages(team_id="t1", channel_id="ch-1", top=10)
    mock_list.assert_called_once_with(mock_client_cls.return_value, "t1", "ch-1", top=10)
    assert len(result) == 1


@patch("outpost.config.is_feature_enabled", return_value=True)
@patch("outpost.mcp_server.get_token", return_value="fake-token")
@patch("outpost.api.teams.send_message")
@patch("outpost.mcp_server.GraphClient")
def test_teams_send(mock_client_cls, mock_send, mock_token, mock_feature):
    mock_send.return_value = {"id": "m-new"}
    result = teams_send(team_id="t1", channel_id="ch-1", body="Hello!")
    mock_send.assert_called_once_with(mock_client_cls.return_value, "t1", "ch-1", "Hello!")
    assert result["id"] == "m-new"


@patch("outpost.config.is_feature_enabled", return_value=True)
@patch("outpost.mcp_server.get_token", return_value="fake-token")
@patch("outpost.api.teams.list_files")
@patch("outpost.api.teams.get_channel_files_folder")
@patch("outpost.mcp_server.GraphClient")
def test_teams_files(mock_client_cls, mock_folder, mock_list, mock_token, mock_feature):
    mock_folder.return_value = {
        "id": "folder-1",
        "parentReference": {"driveId": "drive-1"},
    }
    mock_list.return_value = [{"id": "f1", "name": "doc.pdf"}]
    result = teams_files(team_id="t1", channel_id="ch-1")
    mock_folder.assert_called_once()
    mock_list.assert_called_once_with(mock_client_cls.return_value, "drive-1", "folder-1")
    assert len(result) == 1


@patch("outpost.config.is_feature_enabled", return_value=True)
@patch("outpost.mcp_server.get_token", return_value="fake-token")
@patch("outpost.api.teams.download_file")
@patch("outpost.mcp_server.GraphClient")
def test_teams_download(mock_client_cls, mock_download, mock_token, mock_feature, tmp_path):
    mock_download.return_value = ("report.pdf", b"content")
    with patch("outpost.mcp_server.get_workspace_dir", return_value=tmp_path):
        result = teams_download(drive_id="drive-1", item_id="file-1")
    mock_download.assert_called_once_with(mock_client_cls.return_value, "drive-1", "file-1")
    assert result["filename"] == "report.pdf"
    assert result["workspace_path"] == str(tmp_path / "report.pdf")
    assert result["size"] == 7
    assert (tmp_path / "report.pdf").read_bytes() == b"content"


@patch("outpost.config.is_feature_enabled", return_value=True)
@patch("outpost.mcp_server.get_token", return_value="fake-token")
@patch("outpost.api.teams.upload_file")
@patch("outpost.mcp_server.GraphClient")
def test_teams_upload(mock_client_cls, mock_upload, mock_token, mock_feature, tmp_path):
    # Create a file in the workspace
    (tmp_path / "modified.txt").write_text("new content")
    mock_upload.return_value = {"id": "item-1", "name": "modified.txt"}
    with patch("outpost.mcp_server.get_workspace_dir", return_value=tmp_path):
        result = teams_upload(drive_id="drive-1", parent_item_id="folder-1", filename="modified.txt")
    mock_upload.assert_called_once_with(
        mock_client_cls.return_value, "drive-1", "folder-1", "modified.txt", b"new content"
    )
    assert result["name"] == "modified.txt"


@patch("outpost.config.is_feature_enabled", return_value=True)
def test_teams_upload_file_not_found(mock_feature, tmp_path):
    with patch("outpost.mcp_server.get_workspace_dir", return_value=tmp_path):
        with pytest.raises(RuntimeError, match="File not found"):
            teams_upload(drive_id="drive-1", parent_item_id="folder-1", filename="nonexistent.txt")


def test_teams_workspace_list(tmp_path):
    (tmp_path / "file1.txt").write_text("hello")
    (tmp_path / "file2.pdf").write_bytes(b"pdf content")
    with patch("outpost.mcp_server.get_workspace_dir", return_value=tmp_path):
        result = teams_workspace_list()
    assert len(result) == 2
    names = {f["filename"] for f in result}
    assert "file1.txt" in names
    assert "file2.pdf" in names


def test_teams_workspace_list_empty(tmp_path):
    with patch("outpost.mcp_server.get_workspace_dir", return_value=tmp_path):
        result = teams_workspace_list()
    assert result == []


def test_teams_workspace_read(tmp_path):
    (tmp_path / "notes.txt").write_text("hello world", encoding="utf-8")
    with patch("outpost.mcp_server.get_workspace_dir", return_value=tmp_path):
        result = teams_workspace_read(filename="notes.txt")
    assert result["filename"] == "notes.txt"
    assert result["content"] == "hello world"


def test_teams_workspace_read_binary(tmp_path):
    (tmp_path / "report.docx").write_bytes(b"\x50\x4b\x03\x04\x00\x00\xff\xff")
    with patch("outpost.mcp_server.get_workspace_dir", return_value=tmp_path):
        result = teams_workspace_read(filename="report.docx")
    assert result["binary"] is True
    assert result["filename"] == "report.docx"
    assert result["size"] == 8
    assert "filesystem server" in result["hint"]
    assert "content" not in result


def test_teams_workspace_read_not_found(tmp_path):
    with patch("outpost.mcp_server.get_workspace_dir", return_value=tmp_path):
        with pytest.raises(RuntimeError, match="File not found"):
            teams_workspace_read(filename="missing.txt")


def test_teams_workspace_write(tmp_path):
    with patch("outpost.mcp_server.get_workspace_dir", return_value=tmp_path):
        result = teams_workspace_write(filename="output.txt", content="new content")
    assert result["filename"] == "output.txt"
    assert (tmp_path / "output.txt").read_text(encoding="utf-8") == "new content"


def test_teams_workspace_write_overwrite(tmp_path):
    (tmp_path / "existing.txt").write_text("old content")
    with patch("outpost.mcp_server.get_workspace_dir", return_value=tmp_path):
        teams_workspace_write(filename="existing.txt", content="updated")
    assert (tmp_path / "existing.txt").read_text(encoding="utf-8") == "updated"


# ── teams_workspace_extract ────────────────────────────────────────────


def test_teams_workspace_extract_docx(tmp_path):
    from docx import Document
    doc = Document()
    doc.add_paragraph("Hello from Outpost")
    doc.add_paragraph("Second paragraph")
    doc.save(str(tmp_path / "test.docx"))
    with patch("outpost.mcp_server.get_workspace_dir", return_value=tmp_path):
        result = teams_workspace_extract(filename="test.docx")
    assert result["format"] == ".docx"
    assert "Hello from Outpost" in result["content"]
    assert "Second paragraph" in result["content"]


def test_teams_workspace_extract_xlsx(tmp_path):
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    ws.append(["Name", "Value"])
    ws.append(["Alpha", 42])
    wb.save(str(tmp_path / "test.xlsx"))
    with patch("outpost.mcp_server.get_workspace_dir", return_value=tmp_path):
        result = teams_workspace_extract(filename="test.xlsx")
    assert result["format"] == ".xlsx"
    assert "Alpha" in result["content"]
    assert "42" in result["content"]


def test_teams_workspace_extract_pptx(tmp_path):
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Slide Title"
    slide.placeholders[1].text = "Bullet point"
    prs.save(str(tmp_path / "test.pptx"))
    with patch("outpost.mcp_server.get_workspace_dir", return_value=tmp_path):
        result = teams_workspace_extract(filename="test.pptx")
    assert result["format"] == ".pptx"
    assert "Slide Title" in result["content"]
    assert "Bullet point" in result["content"]


def test_teams_workspace_extract_unsupported(tmp_path):
    (tmp_path / "image.png").write_bytes(b"\x89PNG")
    with patch("outpost.mcp_server.get_workspace_dir", return_value=tmp_path):
        with pytest.raises(RuntimeError, match="Unsupported file type"):
            teams_workspace_extract(filename="image.png")


def test_teams_workspace_extract_not_found(tmp_path):
    with patch("outpost.mcp_server.get_workspace_dir", return_value=tmp_path):
        with pytest.raises(RuntimeError, match="File not found"):
            teams_workspace_extract(filename="missing.docx")
