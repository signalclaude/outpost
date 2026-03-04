"""MCP server exposing Outpost operations as tools for Claude Desktop."""

import base64
from datetime import datetime, timedelta
from typing import Optional

from fastmcp import FastMCP

from outpost.api.client import GraphClient
from outpost.auth import get_auth_status, get_token
from outpost.config import clean_workspace, get_workspace_dir, load_config
from outpost.utils.dates import (
    parse_natural_date,
    parse_natural_datetime,
    to_graph_date,
    today_range,
)

mcp = FastMCP(
    "outpost",
    instructions="Manage Microsoft Outlook tasks, calendar, email, contacts, and Teams via Graph API.",
)


def create_mcp_server(auth=None):
    """Create a new FastMCP server instance, optionally with auth.

    For stdio transport, returns the module-level `mcp` (no auth needed).
    For network transports (SSE, streamable-http), pass an auth provider.
    """
    if auth is None:
        return mcp
    server = FastMCP(
        "outpost",
        instructions="Manage Microsoft Outlook tasks, calendar, email, contacts, and Teams via Graph API.",
        auth=auth,
    )
    server.mount(mcp)
    return server

# Clean workspace on MCP server startup
clean_workspace()


def _get_client() -> GraphClient:
    """Get an authenticated GraphClient or raise RuntimeError."""
    token = get_token()
    if token is None:
        raise RuntimeError(
            "Not authenticated. Run 'outpost setup' in a terminal first to connect your Microsoft account."
        )
    return GraphClient(token=token)


@mcp.tool()
def task_add(
    title: str,
    due: Optional[str] = None,
    list_name: Optional[str] = None,
    priority: str = "normal",
) -> dict:
    """Add a new task to Microsoft To Do.

    Args:
        title: Task title
        due: Due date — natural language ('tomorrow', 'next friday') or ISO ('2026-03-15')
        list_name: Task list name (uses default list if omitted)
        priority: low, normal, or high
    """
    from outpost.api.tasks import add_task

    client = _get_client()
    due_date = to_graph_date(parse_natural_date(due)) if due else None
    return add_task(client, title, due_date=due_date, list_name=list_name, priority=priority)


@mcp.tool()
def task_list(
    due: Optional[str] = None,
    list_name: Optional[str] = None,
) -> list[dict]:
    """List tasks from Microsoft To Do.

    Args:
        due: Filter by due date — natural language or ISO format
        list_name: Task list name (uses default list if omitted)
    """
    from outpost.api.tasks import list_tasks

    client = _get_client()
    due_filter = to_graph_date(parse_natural_date(due)) if due else None
    return list_tasks(client, due_filter=due_filter, list_name=list_name)


@mcp.tool()
def cal_add(
    title: str,
    start: str,
    end: Optional[str] = None,
    duration: Optional[int] = None,
    attendees: Optional[str] = None,
    show_as: Optional[str] = None,
    timezone: Optional[str] = None,
) -> dict:
    """Add a calendar event to Outlook.

    Args:
        title: Event title/subject
        start: Start time — natural language ('tomorrow 9am') or ISO ('2026-03-15T09:00')
        end: End time (optional if duration is given)
        duration: Duration in minutes (default 30 if neither end nor duration given)
        attendees: Comma-separated attendee email addresses (optional)
        show_as: Free/busy status — free, tentative, busy, oof (Out of Office), or workingElsewhere (optional)
        timezone: IANA timezone (e.g. 'America/New_York'). Defaults to value from outpost config.
    """
    from outpost.api.calendar import create_event

    client = _get_client()
    tz = timezone or load_config().get("timezone", "UTC")
    start_dt = parse_natural_datetime(start)
    end_dt = parse_natural_datetime(end) if end else None
    attendee_list = [a.strip() for a in attendees.split(",")] if attendees else None
    return create_event(client, title, start_dt, end=end_dt, duration=duration, attendees=attendee_list, show_as=show_as, timezone=tz)


@mcp.tool()
def cal_today() -> list[dict]:
    """Get today's calendar events from Outlook."""
    from outpost.api.calendar import get_today_events

    client = _get_client()
    return get_today_events(client)


@mcp.tool()
def cal_list(
    date: Optional[str] = None,
    week: bool = False,
) -> list[dict]:
    """List calendar events from Outlook.

    Args:
        date: Show events for a specific date (natural language or ISO)
        week: If true, show this week's events (default if no date given)
    """
    from outpost.api.calendar import get_calendar_view, get_week_events

    client = _get_client()

    if date:
        d = parse_natural_date(date)
        start = datetime(d.year, d.month, d.day, 0, 0, 0)
        end = start + timedelta(days=1)
        return get_calendar_view(client, start, end)
    else:
        return get_week_events(client)


@mcp.tool()
def cal_next(count: int = 1) -> list[dict]:
    """Get the next upcoming calendar event(s) from Outlook.

    Args:
        count: Number of upcoming events to return (default 1)
    """
    from outpost.api.calendar import get_next_events

    client = _get_client()
    return get_next_events(client, count=count)


@mcp.tool()
def task_update(
    task_id: str,
    title: Optional[str] = None,
    due: Optional[str] = None,
    priority: Optional[str] = None,
    list_name: Optional[str] = None,
) -> dict:
    """Update an existing task in Microsoft To Do.

    Args:
        task_id: Task ID to update
        title: New title (optional)
        due: New due date — natural language or ISO format (optional)
        priority: New priority: low, normal, or high (optional)
        list_name: Task list name (uses default list if omitted)
    """
    from outpost.api.tasks import update_task

    client = _get_client()
    due_date = to_graph_date(parse_natural_date(due)) if due else None
    return update_task(client, task_id, list_name=list_name, title=title, due_date=due_date, priority=priority)


@mcp.tool()
def task_complete(task_id: str, list_name: Optional[str] = None) -> dict:
    """Mark a task as completed in Microsoft To Do.

    Args:
        task_id: Task ID to complete
        list_name: Task list name (uses default list if omitted)
    """
    from outpost.api.tasks import complete_task

    client = _get_client()
    return complete_task(client, task_id, list_name=list_name)


@mcp.tool()
def task_delete(task_id: str, list_name: Optional[str] = None) -> dict:
    """Delete a task from Microsoft To Do.

    Args:
        task_id: Task ID to delete
        list_name: Task list name (uses default list if omitted)
    """
    from outpost.api.tasks import delete_task

    client = _get_client()
    return delete_task(client, task_id, list_name=list_name)


@mcp.tool()
def task_lists() -> list[dict]:
    """List all task lists in Microsoft To Do."""
    from outpost.api.tasks import get_task_lists

    client = _get_client()
    return get_task_lists(client)


@mcp.tool()
def task_lists_create(name: str) -> dict:
    """Create a new task list in Microsoft To Do.

    Args:
        name: Display name for the new task list
    """
    from outpost.api.tasks import create_task_list

    client = _get_client()
    return create_task_list(client, name)


@mcp.tool()
def task_lists_delete(list_id: str) -> dict:
    """Delete a task list from Microsoft To Do.

    Args:
        list_id: Task list ID to delete
    """
    from outpost.api.tasks import delete_task_list

    client = _get_client()
    return delete_task_list(client, list_id)


@mcp.tool()
def cal_update(
    event_id: str,
    title: Optional[str] = None,
    start: Optional[str] = None,
    end: Optional[str] = None,
    attendees: Optional[str] = None,
    show_as: Optional[str] = None,
    timezone: Optional[str] = None,
) -> dict:
    """Update an existing calendar event in Outlook.

    Args:
        event_id: Event ID to update
        title: New title (optional)
        start: New start time — natural language or ISO format (optional)
        end: New end time — natural language or ISO format (optional)
        attendees: Comma-separated attendee emails (optional)
        show_as: Free/busy status — free, tentative, busy, oof (Out of Office), or workingElsewhere (optional)
        timezone: IANA timezone (e.g. 'America/New_York'). Defaults to value from outpost config.
    """
    from outpost.api.calendar import update_event

    client = _get_client()
    tz = timezone or load_config().get("timezone", "UTC")
    start_dt = parse_natural_datetime(start) if start else None
    end_dt = parse_natural_datetime(end) if end else None
    attendee_list = [a.strip() for a in attendees.split(",")] if attendees else None
    return update_event(client, event_id, title=title, start=start_dt, end=end_dt, attendees=attendee_list, show_as=show_as, timezone=tz)


@mcp.tool()
def cal_delete(event_id: str) -> dict:
    """Delete a calendar event from Outlook.

    Args:
        event_id: Event ID to delete
    """
    from outpost.api.calendar import delete_event

    client = _get_client()
    return delete_event(client, event_id)


@mcp.tool()
def mail_list(folder: str = "inbox", top: int = 25, unread: bool = False) -> list[dict]:
    """List email messages from Outlook.

    Args:
        folder: Mail folder (inbox, sentitems, drafts, deleteditems)
        top: Maximum number of messages to return
        unread: If true, show only unread messages
    """
    from outpost.api.mail import list_messages

    client = _get_client()
    return list_messages(client, folder=folder, top=top, unread=unread)


@mcp.tool()
def mail_read(message_id: str) -> dict:
    """Read a specific email message from Outlook.

    Args:
        message_id: Message ID to read
    """
    from outpost.api.mail import get_message

    client = _get_client()
    return get_message(client, message_id)


@mcp.tool()
def mail_send(to: str, subject: str, body: str, cc: Optional[str] = None) -> dict:
    """Send an email via Outlook.

    Args:
        to: Recipient email address (comma-separated for multiple)
        subject: Email subject
        body: Email body text
        cc: CC recipients (comma-separated, optional)
    """
    from outpost.api.mail import send_message

    client = _get_client()
    to_list = [addr.strip() for addr in to.split(",")]
    cc_list = [addr.strip() for addr in cc.split(",")] if cc else None
    return send_message(client, to_list, subject, body, cc=cc_list)


@mcp.tool()
def mail_reply(message_id: str, body: str) -> dict:
    """Reply to an email in Outlook.

    Args:
        message_id: Message ID to reply to
        body: Reply text
    """
    from outpost.api.mail import reply_message

    client = _get_client()
    return reply_message(client, message_id, body)


@mcp.tool()
def mail_delete(message_id: str) -> dict:
    """Delete an email from Outlook.

    Args:
        message_id: Message ID to delete
    """
    from outpost.api.mail import delete_message

    client = _get_client()
    return delete_message(client, message_id)


@mcp.tool()
def mail_search(query: str, top: int = 25) -> list[dict]:
    """Search email messages across all folders in Outlook.

    Args:
        query: Search query string
        top: Maximum number of results
    """
    from outpost.api.mail import search_messages

    client = _get_client()
    return search_messages(client, query, top=top)


@mcp.tool()
def mail_attachments(message_id: str) -> list[dict]:
    """List attachments for an email message.

    Args:
        message_id: Message ID to list attachments for
    """
    from outpost.api.mail import list_attachments

    client = _get_client()
    return list_attachments(client, message_id)


@mcp.tool()
def mail_download_attachment(message_id: str, attachment_id: str) -> dict:
    """Download an email attachment. Returns filename and base64-encoded content.

    Args:
        message_id: Message ID
        attachment_id: Attachment ID to download
    """
    from outpost.api.mail import download_attachment

    client = _get_client()
    filename, content = download_attachment(client, message_id, attachment_id)
    return {"filename": filename, "content_base64": base64.b64encode(content).decode("ascii")}


@mcp.tool()
def contact_list(top: int = 50) -> list[dict]:
    """List contacts from Outlook.

    Args:
        top: Maximum number of contacts to return
    """
    from outpost.api.contacts import list_contacts

    client = _get_client()
    return list_contacts(client, top=top)


@mcp.tool()
def contact_search(query: str, top: int = 25) -> list[dict]:
    """Search contacts in Outlook by name or email.

    Args:
        query: Search query (name or email)
        top: Maximum number of results
    """
    from outpost.api.contacts import search_contacts

    client = _get_client()
    return search_contacts(client, query, top=top)


@mcp.tool()
def auth_status() -> dict:
    """Check if the user is authenticated with Microsoft."""
    return get_auth_status()


# ── Teams tools (feature-gated) ─────────────────────────────────────────────


def _require_teams():
    """Raise RuntimeError if the teams feature is not enabled."""
    from outpost.config import is_feature_enabled

    if not is_feature_enabled("teams"):
        raise RuntimeError(
            "The 'teams' feature is not enabled. Run 'outpost setup' and enable it first."
        )


@mcp.tool()
def teams_list() -> list[dict]:
    """List your joined Microsoft Teams. Requires the Teams feature to be enabled in outpost setup."""
    _require_teams()
    from outpost.api.teams import list_teams

    client = _get_client()
    return list_teams(client)


@mcp.tool()
def teams_channels(team_id: str) -> list[dict]:
    """List channels in a Microsoft Team.

    Args:
        team_id: Team ID
    """
    _require_teams()
    from outpost.api.teams import list_channels

    client = _get_client()
    return list_channels(client, team_id)


@mcp.tool()
def teams_messages(team_id: str, channel_id: str, top: int = 25) -> list[dict]:
    """Read messages from a Teams channel.

    Args:
        team_id: Team ID
        channel_id: Channel ID
        top: Maximum number of messages to return
    """
    _require_teams()
    from outpost.api.teams import list_messages

    client = _get_client()
    return list_messages(client, team_id, channel_id, top=top)


@mcp.tool()
def teams_send(team_id: str, channel_id: str, body: str) -> dict:
    """Send a message to a Teams channel.

    Args:
        team_id: Team ID
        channel_id: Channel ID
        body: Message content
    """
    _require_teams()
    from outpost.api.teams import send_message

    client = _get_client()
    return send_message(client, team_id, channel_id, body)


@mcp.tool()
def teams_files(team_id: str, channel_id: str) -> list[dict]:
    """List files in a Teams channel's SharePoint folder.

    Args:
        team_id: Team ID
        channel_id: Channel ID
    """
    _require_teams()
    from outpost.api.teams import get_channel_files_folder, list_files

    client = _get_client()
    folder = get_channel_files_folder(client, team_id, channel_id)
    drive_id = folder["parentReference"]["driveId"]
    item_id = folder["id"]
    return list_files(client, drive_id, item_id)


@mcp.tool()
def teams_chats(top: int = 25) -> list[dict]:
    """List your recent Teams chats (1:1, group, and meeting). Requires the Teams feature.

    Args:
        top: Maximum number of chats to return
    """
    _require_teams()
    from outpost.api.teams import list_chats

    client = _get_client()
    return list_chats(client, top=top)


@mcp.tool()
def teams_chat_messages(chat_id: str, top: int = 25) -> list[dict]:
    """Read messages from a Teams chat (1:1 or group).

    Args:
        chat_id: Chat ID (from teams_chats)
        top: Maximum number of messages to return
    """
    _require_teams()
    from outpost.api.teams import list_chat_messages

    client = _get_client()
    return list_chat_messages(client, chat_id, top=top)


@mcp.tool()
def teams_chat_send(chat_id: str, body: str) -> dict:
    """Send a message to a Teams chat (1:1 or group).

    Args:
        chat_id: Chat ID (from teams_chats)
        body: Message content
    """
    _require_teams()
    from outpost.api.teams import send_chat_message

    client = _get_client()
    return send_chat_message(client, chat_id, body)


@mcp.tool()
def teams_download(drive_id: str, item_id: str) -> dict:
    """Download a file from Teams/SharePoint to the local workspace.

    The file is saved to the local workspace. To read or analyze the file, use the
    filesystem server to read it at the returned workspace_path. Use teams_workspace_read
    only for plain text files. Use teams_upload to push changes back.

    Args:
        drive_id: Drive ID
        item_id: Item ID
    """
    _require_teams()
    from outpost.api.teams import download_file

    client = _get_client()
    filename, content = download_file(client, drive_id, item_id)
    workspace = get_workspace_dir()
    filepath = workspace / filename
    filepath.write_bytes(content)
    return {"filename": filename, "workspace_path": str(filepath), "size": len(content)}


@mcp.tool()
def teams_upload(drive_id: str, parent_item_id: str, filename: str) -> dict:
    """Upload a file from the workspace to Teams/SharePoint.

    Best suited for text-based files created with teams_workspace_write (markdown,
    CSV, plain text). Binary files (DOCX, XLSX, etc.) should be treated as read-only
    — use teams_workspace_extract to read them, not modify and re-upload.

    Args:
        drive_id: Drive ID
        parent_item_id: Parent folder item ID (from teams_files)
        filename: Name of the file in the workspace to upload
    """
    _require_teams()
    from outpost.api.teams import upload_file

    workspace = get_workspace_dir()
    filepath = workspace / filename
    if not filepath.exists():
        raise RuntimeError(f"File not found in workspace: {filename}")
    content = filepath.read_bytes()
    client = _get_client()
    return upload_file(client, drive_id, parent_item_id, filename, content)


@mcp.tool()
def teams_workspace_list() -> list[dict]:
    """List files currently in the transient workspace directory.

    For binary files (docx, pdf, xlsx, etc.), read them directly at the returned
    path using the filesystem server.
    """
    workspace = get_workspace_dir()
    files = []
    for item in workspace.iterdir():
        if item.is_file():
            files.append({
                "filename": item.name,
                "size": item.stat().st_size,
                "path": str(item),
            })
    return files


@mcp.tool()
def teams_workspace_read(filename: str) -> dict:
    """Read a text file from the workspace directory.

    For binary files (docx, pdf, xlsx, etc.), this tool returns metadata and a path
    hint instead of content. Use the filesystem server to read binary files directly.

    Args:
        filename: Name of the file to read
    """
    workspace = get_workspace_dir()
    filepath = workspace / filename
    if not filepath.exists():
        raise RuntimeError(f"File not found in workspace: {filename}")
    try:
        content = filepath.read_text(encoding="utf-8")
    except (UnicodeDecodeError, ValueError):
        return {
            "filename": filename,
            "size": filepath.stat().st_size,
            "path": str(filepath),
            "binary": True,
            "hint": f"Binary file — read it directly at {filepath} using the filesystem server",
        }
    return {"filename": filename, "content": content}


@mcp.tool()
def teams_workspace_write(filename: str, content: str) -> dict:
    """Write a text file to the workspace directory for sharing via Teams.

    For text-based formats only (markdown, CSV, plain text, etc.). Do not use this
    to create or overwrite binary files (DOCX, XLSX, PPTX, PDF) — use
    teams_workspace_extract to read those instead.

    Args:
        filename: Name of the file to write (e.g. 'notes.md', 'data.csv')
        content: Text content to write
    """
    workspace = get_workspace_dir()
    filepath = workspace / filename
    filepath.write_text(content, encoding="utf-8")
    return {"filename": filename, "size": filepath.stat().st_size, "path": str(filepath)}


EXTRACTORS = {
    ".docx": "_extract_docx",
    ".xlsx": "_extract_xlsx",
    ".pptx": "_extract_pptx",
    ".pdf": "_extract_pdf",
}


def _extract_docx(filepath) -> str:
    from docx import Document
    doc = Document(str(filepath))
    return "\n".join(p.text for p in doc.paragraphs)


def _extract_xlsx(filepath) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(str(filepath), read_only=True, data_only=True)
    lines = []
    for sheet in wb.sheetnames:
        lines.append(f"--- Sheet: {sheet} ---")
        for row in wb[sheet].iter_rows(values_only=True):
            lines.append("\t".join(str(c) if c is not None else "" for c in row))
    wb.close()
    return "\n".join(lines)


def _extract_pptx(filepath) -> str:
    from pptx import Presentation
    prs = Presentation(str(filepath))
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                lines.append(shape.text_frame.text)
    return "\n".join(lines)


def _extract_pdf(filepath) -> str:
    import pdfplumber
    with pdfplumber.open(str(filepath)) as pdf:
        return "\n".join(page.extract_text() or "" for page in pdf.pages)


@mcp.tool()
def teams_workspace_extract(filename: str) -> dict:
    """Extract text content from a binary file in the workspace.

    Supports DOCX, XLSX, PPTX, and PDF files. Returns the extracted plain text
    content. Much faster than reading binary files through the filesystem server.

    Args:
        filename: Name of the file in the workspace to extract text from
    """
    workspace = get_workspace_dir()
    filepath = workspace / filename
    if not filepath.exists():
        raise RuntimeError(f"File not found in workspace: {filename}")

    ext = filepath.suffix.lower()
    extractor_name = EXTRACTORS.get(ext)
    if not extractor_name:
        supported = ", ".join(EXTRACTORS.keys())
        raise RuntimeError(f"Unsupported file type: {ext}. Supported: {supported}")

    extractor = globals()[extractor_name]
    text = extractor(filepath)
    return {"filename": filename, "format": ext, "content": text}


def main():
    """Run the MCP server over stdio."""
    mcp.run()


if __name__ == "__main__":
    main()
